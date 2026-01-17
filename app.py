import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests, json, io

# 1. ПОРЯДОК СТИЛЕЙ
THEMES = {
    "SCHOOL STYLE": {"acc": (50, 150, 50), "icon": "✏️", "left": 1.5, "width": 10.3, "dark": True},
    "GIRLY STYLE": {"acc": (255, 105, 180), "icon": "🌸", "left": 1.5, "width": 10.3, "dark": False},
    "MODERN GRADIENT": {"acc": (0, 102, 204), "icon": "➔", "left": 1.0, "width": 11.3, "dark": False},
    "MINIMALIST": {"acc": (100, 100, 100), "icon": "◈", "left": 1.5, "width": 10.3, "dark": False},
    "NEON NIGHT": {"acc": (0, 255, 150), "icon": "⚡", "left": 1.0, "width": 11.3, "dark": True},
    "BUSINESS PRO": {"acc": (0, 80, 180), "icon": "✔", "left": 1.0, "width": 11.3, "dark": False},
    "SUNSET STYLE": {"acc": (255, 140, 0), "icon": "☀️", "left": 1.0, "width": 11.3, "dark": True},
    "LUFFY STYLE": {"acc": (200, 30, 30), "icon": "⚓", "left": 5.8, "width": 7.0, "dark": False},
}

AI_KEY = st.secrets.get("GROQ_API_KEY", "")

def ask_ai(topic, slides, lang):
    if not AI_KEY: return None
    # Усиленная инструкция для таджикского языка
    sys_content = (f"You are a professional professor. Write in {lang}. "
                   "For Tajik language, use only pure literary Tajik grammar (Забони адабии тоҷикӣ). "
                   "Output ONLY valid JSON.")
    
    prompt = (f"Create presentation '{topic}'. Slides: {slides}. "
              "Each slide 'intro' MUST be 80-160 words. "
              "Create 10 quiz questions. Return JSON: {'slides': [{'title': '...', 'intro': '...'}], "
              "'quiz': [{'q': '...', 'a': 'A', 'o': ['A-..','B-..','C-..']}]}")
    try:
        r = requests.post("https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {AI_KEY}"},
            json={"model": "llama-3.3-70b-versatile", "messages": [
                {"role": "system", "content": sys_content},
                {"role": "user", "content": prompt}
            ], "response_format": {"type": "json_object"}}, timeout=120)
        return json.loads(r.json()["choices"][0]["message"]["content"])
    except: return None

def make_pptx(data, style_name, font_size):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    theme = THEMES[style_name]
    slides_data = data.get('slides', [])
    
    for s in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txt_rgb = RGBColor(255, 255, 255) if theme["dark"] else RGBColor(30, 30, 30)
        
        try: slide.shapes.add_picture(f"{style_name}.jpg", 0, 0, width=prs.slide_width, height=prs.slide_height)
        except: pass
        
        # ЗАГОЛОВОК
        p_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0)).text_frame.paragraphs[0]
        p_t.text = f"{theme['icon']} {str(s.get('title', '')).upper()}"
        p_t.font.name, p_t.font.size, p_t.font.bold = 'Times New Roman', Pt(40), True
        p_t.font.color.rgb = RGBColor(*theme["acc"])
        
        # ТЕКСТ (С выбранным размером)
        tf = slide.shapes.add_textbox(Inches(theme["left"]), Inches(1.5), Inches(theme["width"]), Inches(5.5)).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = str(s.get('intro', ''))
        p.font.name, p.font.size, p.font.color.rgb = 'Times New Roman', Pt(font_size), txt_rgb
        tf.line_spacing = 1.15
        
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf

st.set_page_config(page_title="SLIDEX PRO", layout="wide")
# Вместо старого заголовка ставим твой прямоугольный логотип
st.image("Logo.jpg", use_container_width=True)

if "data" not in st.session_state: st.session_state.data = None
if "test_key" not in st.session_state: st.session_state.test_key = 0
if "submitted" not in st.session_state: st.session_state.submitted = False

with st.sidebar:
    # 1. Делаем логотип маленьким и КЛИКАБЕЛЬНЫМ (всё в одном)
    import base64
    def get_base64(file_path):
        with open(file_path, "rb") as f:
            return base64.b64encode(f.read()).decode()

    try:
        # Используем твой файл 1000021955.jpg
        img_data = get_base64("1000021955.jpg")
        st.markdown(
            f"""
            <div style="text-align: center;">
                <a href="https://amin-cloud-copy-8f1d0b41.base44.app/" target="_blank" style="text-decoration: none;">
                    <img src="data:image/png;base64,{img_data}" width="120" style="border-radius: 10px;">
                    <p style="color: #00d4ff; font-size: 12px; margin-top: 5px;">Перейти в AminCloud</p>
                </a>
            </div>
            """,
            unsafe_allow_html=True
        )
    except:
        st.link_button("🌐 AminCloud", "https://amin-cloud-copy-8f1d0b41.base44.app/")

    st.divider()

    st.link_button("🌐 Перейти в AminCloud", "https://amin-cloud-copy-8f1d0b41.base44.app/", use_container_width=True)
    st.divider()
    t_input = st.text_input("Тема презентации")
    s_count = st.slider("Слайды (от 2 до 12)", 2, 12, 6)
    f_size = st.slider("Размер шрифта в файле", 26, 40, 32)
    lang_choice = st.selectbox("Язык", ["Russian", "Tajik", "English"])
    style_sel = st.selectbox("Стиль", list(THEMES.keys()))
    user_code = st.text_input("Код доступа", type="password") 
    
    if st.button("🚀 Сгенерировать"):
        res = ask_ai(t_input, s_count, lang_choice)
        if res:
            st.session_state.data = res
            st.session_state.test_key += 1
            st.session_state.submitted = False
            st.rerun()

if st.session_state.data:
    st.header(f"Просмотр контента")
    for i, s in enumerate(st.session_state.data['slides']):
        with st.expander(f"Слайд {i+1}: {s.get('title')}"):
            st.write(s.get('intro'))

    # ЛОГИКА ДОСТУПА
    if user_code == "SX-369":
        st.success("🔓 Режим активен")
        st.download_button("📥 СКАЧАТЬ PPTX", make_pptx(st.session_state.data, style_sel, f_size), f"{t_input}.pptx")
    else:
        st.header("✅ Проверка знаний (Тест)")
        quiz = st.session_state.data.get('quiz', [])[:10]
        user_ans = []
        for i, q in enumerate(quiz):
            ans = st.selectbox(f"Вопрос {i+1}: {q['q']}", ["-- Выберите --"] + q['o'], key=f"q_{i}_{st.session_state.test_key}")
            user_ans.append(ans)

        if st.button("Проверить результат"):
            if "-- Выберите --" in user_ans: st.warning("Ответьте на все вопросы!")
            else:
                score = 0
                for i, q in enumerate(quiz):
                    if user_ans[i][0] == q['a']: score += 1
                
                st.subheader(f"Ваш результат: {score}/10")
                if score >= 8:
                    st.balloons()
                    st.download_button("📥 СКАЧАТЬ PPTX", make_pptx(st.session_state.data, style_sel, f_size), f"{t_input}.pptx")
                else:
                    st.error("Нужно минимум 8 правильных ответов.")
