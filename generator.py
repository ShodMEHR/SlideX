# generator.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
from styles import ICONS

def clamp_intro(text, min_w=80, max_w=160):
    """Обеспечивает объем текста от 80 до 160 слов"""
    words = text.split()
    if len(words) < min_w:
        words += words[: (min_w - len(words))]
    return " ".join(words[:max_w])

def make_pptx(data, theme, style):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5) # Формат 16:9
    current_icon = ICONS.get(style, "• ")

    for s in data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txt_rgb = RGBColor(30, 30, 30) # По умолчанию темный текст
        acc_rgb = RGBColor(*theme["acc"])
        l_m, w_m = Inches(1.0), Inches(11.3) # Стандартные отступы

        # Пытаемся наложить картинку фона
        try:
            slide.shapes.add_picture(f"{style}.jpg", 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            # Специфическая верстка под картинки
            if style == "LUFFY STYLE": 
                l_m, w_m = Inches(5.5), Inches(7.3) # Текст справа
            elif style == "GIRLY STYLE":
                # Белая подложка для читаемости
                rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.3), Inches(5.8))
                rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(255, 255, 255); rect.fill.alpha = 0.8
                l_m, w_m = Inches(2.0), Inches(9.3)
            elif style in ["SCHOOL STYLE", "NEON NIGHT", "SUNSET STYLE"]: 
                txt_rgb = RGBColor(255, 255, 255) # Белый текст для темных фонов
        except:
            # Если картинки нет, просто заливаем серым
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(240, 240, 240)

        # Создаем Заголовок
        tb_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
        p_t = tb_t.text_frame.paragraphs[0]
        p_t.text = str(s.get("title","")).upper()
        p_t.font.size, p_t.font.bold, p_t.font.color.rgb = Pt(36), True, acc_rgb

        # Создаем блок текста (80-160 слов)
        box = slide.shapes.add_textbox(l_m, Inches(1.4), w_m, Inches(5.5))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = clamp_intro(str(s.get("intro","")))
        p.font.size, p.font.color.rgb = Pt(14), txt_rgb

        # Добавляем пункты с иконками
        for pt in s.get("points", []):
            pp = tf.add_paragraph()
            pp.text = f"{current_icon}{pt}"
            pp.font.size, pp.font.bold, pp.font.color.rgb = Pt(12), True, acc_rgb

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
